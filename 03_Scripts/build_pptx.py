"""Build the RoadSoS hackathon presentation deck (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import os
OUT = os.path.dirname(os.path.abspath(__file__)) if '__file__' in locals() else os.getcwd()
NAVY = RGBColor(0x0f, 0x34, 0x60)
NAVY_DARK = RGBColor(0x08, 0x1a, 0x36)
AMBER = RGBColor(0xf5, 0xa6, 0x23)
RED = RGBColor(0xe6, 0x39, 0x46)
TEAL = RGBColor(0x2a, 0x9d, 0x8f)
PURPLE = RGBColor(0x7c, 0x3a, 0xed)
WHITE = RGBColor(0xff, 0xff, 0xff)
GREY = RGBColor(0x47, 0x55, 0x69)
LIGHT_GREY = RGBColor(0x9c, 0xa3, 0xaf)
BG = RGBColor(0xf8, 0xfa, 0xfc)
CARD = RGBColor(0xff, 0xff, 0xff)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, italic=False,
             color=NAVY, align='left', font='Calibri', anchor='top'):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if anchor == 'middle':
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == 'center': p.alignment = PP_ALIGN.CENTER
    elif align == 'right': p.alignment = PP_ALIGN.RIGHT
    elif align == 'left': p.alignment = PP_ALIGN.LEFT
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            run = p.add_run()
        else:
            new_p = tf.add_paragraph()
            new_p.alignment = p.alignment
            run = new_p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill_color, line=None, corner=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_title_bar(slide, title, subtitle=None):
    # left amber accent
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.4), Inches(0.18), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background(); bar.shadow.inherit = False
    add_text(slide, Inches(0.45), Inches(0.35), Inches(12), Inches(0.6), title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.45), Inches(0.95), Inches(12), Inches(0.4), subtitle, size=14, italic=True, color=GREY)


def add_footer(slide, idx, total=12):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(6), Inches(0.3),
             "RoadSoS  ·  Smart Helmet SOS + Predictive Safety", size=10, color=LIGHT_GREY)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{idx} / {total}", size=10, color=LIGHT_GREY, align='right')


# ============ SLIDE 1 — Title ============
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY_DARK)
# Decorative amber accents
for i, (x, y, r) in enumerate([(11.5, 0.5, 1.2), (0.5, 6.0, 0.8)]):
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r))
    circ.fill.solid(); circ.fill.fore_color.rgb = AMBER if i == 0 else RED
    circ.line.fill.background(); circ.shadow.inherit = False
    # make it lower opacity via XML
    sp = circ.fill.fore_color._xFill
add_text(s, Inches(0.7), Inches(0.6), Inches(8), Inches(0.4),
         "HACKATHON  ·  TWO-WHEELER ROAD SAFETY", size=12, bold=True, color=AMBER)
add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(2.2),
         "RoadSoS", size=96, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.8), Inches(12), Inches(1.0),
         "Predicting the crash. Calling for help.\nWarning the world. Powered by the sun.",
         size=24, italic=True, color=AMBER)
add_text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.5),
         "Smart Helmet SOS + Predictive Safety System",
         size=18, color=WHITE)
add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.4),
         "Shaurya  ·  Hackathon Prototype Blueprint  ·  2026",
         size=12, color=LIGHT_GREY)


# ============ SLIDE 2 — Problem ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "The Problem", "The killer isn't always the crash — it's the wait.")

stats = [("1.5 L", "annual road\ndeaths in India"),
         ("44%", "are two-wheeler\nriders"),
         ("~30 min", "avg ambulance arrival\nsemi-urban India"),
         ("70%", "survival lift if reached\nwithin Golden Hour")]
for i, (val, lbl) in enumerate(stats):
    x = Inches(0.5 + i * 3.1)
    add_rect(s, x, Inches(1.8), Inches(2.9), Inches(2.0), CARD, line=AMBER, corner=True)
    add_text(s, x, Inches(2.0), Inches(2.9), Inches(0.8), val, size=44, bold=True, color=AMBER, align='center')
    add_text(s, x, Inches(2.9), Inches(2.9), Inches(0.9), lbl, size=14, color=NAVY, align='center')

add_text(s, Inches(0.5), Inches(4.1), Inches(12.4), Inches(0.5),
         "Today's reality", size=20, bold=True, color=NAVY)
points = [
    "The rider often loses consciousness — nobody calls for help.",
    "Bystanders don't know the victim's blood group, allergies, or contacts.",
    "Following riders collide into the same wreck — secondary collisions.",
    "Conventional sensors trigger only after impact — too late to prevent it.",
]
for i, p in enumerate(points):
    add_text(s, Inches(0.6), Inches(4.6 + i * 0.42), Inches(0.2), Inches(0.4), "•", size=18, bold=True, color=AMBER)
    add_text(s, Inches(0.95), Inches(4.55 + i * 0.42), Inches(12), Inches(0.4), p, size=15, color=GREY)
add_footer(s, 2)


# ============ SLIDE 3 — Solution ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "Our Solution — RoadSoS", "A predictive, self-aware, mesh-connected helmet powered by the sun.")
features = [
    ("🧠", "PINN + TinyML on edge", "Predicts skid before impact;\nconfirms crash after.", AMBER),
    ("📡", "LoRa rider mesh", "Warns every helmet within\n~3 km — no cell tower needed.", TEAL),
    ("🚑", "Auto-SOS to 108", "Cellular + BLE redundancy.\nMedical profile attached.", RED),
    ("☀", "Solar trickle charge", "Daytime rides keep the\nhelmet topped up — net positive.", AMBER),
    ("📲", "QR + NFC medical access", "Bystanders see first-aid &\nblood group instantly.", PURPLE),
    ("🔦", "Crash beacon mode", "Strobe LEDs + Morse buzzer\nhelp responders find rider.", NAVY),
]
for i, (icon, title, body, color) in enumerate(features):
    col, row = i % 3, i // 3
    x = Inches(0.5 + col * 4.2); y = Inches(1.85 + row * 2.5)
    add_rect(s, x, y, Inches(3.95), Inches(2.25), CARD, line=color, corner=True)
    add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.6), icon, size=30)
    add_text(s, x + Inches(1.0), y + Inches(0.25), Inches(2.9), Inches(0.5), title, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(1.0), Inches(3.7), Inches(1.2), body, size=12.5, color=GREY)
add_footer(s, 3)


# ============ SLIDE 4 — Scenario ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "A Day with Ravi", "Hosur → Bangalore, NH-44. One rainy commute, three save-points.")
scenario = [
    ("06:30", "Boot", "Helmet powers on. Battery 92%. Medical profile encrypted on ESP32-S3."),
    ("07:12", "PINN warns", "Oil patch at km 47. PINN flags skid prob > 0.65. Haptic taps temple. Crisis averted."),
    ("11:40", "Crash", "Truck swerves. IMU >6g + free-fall + tilt >80°. TinyML CNN confirms 'crash'."),
    ("11:40:10", "SOS fires", "108, hospital, Anjali notified. LoRa broadcast goes out. Beacon mode active."),
    ("11:54", "Ambulance", "Golden Hour preserved. Three trailing riders slowed in time. No secondary collision."),
    ("16:00", "Return", "Solar trickle kept battery at 87% after 8 hours."),
]
for i, (t, label, body) in enumerate(scenario):
    y = Inches(1.85 + i * 0.83)
    # time chip
    chip = add_rect(s, Inches(0.5), y, Inches(1.3), Inches(0.65), NAVY, corner=True)
    add_text(s, Inches(0.5), y + Inches(0.13), Inches(1.3), Inches(0.5), t, size=14, bold=True, color=AMBER, align='center')
    # label
    add_text(s, Inches(1.95), y + Inches(0.05), Inches(2.5), Inches(0.6), label, size=18, bold=True, color=NAVY)
    # body
    add_text(s, Inches(4.4), y + Inches(0.12), Inches(8.7), Inches(0.6), body, size=13, color=GREY)
add_footer(s, 4)


# ============ SLIDE 5 — Architecture ============
# Available image area: y from 1.55" to 6.95" (height 5.4") on a 13.333 x 7.5" slide.
# Fit-by-height for each figure based on its viewBox aspect ratio.
def fit_image(slide, path, view_w, view_h, top=Inches(1.55), max_h=Inches(5.4), max_w=Inches(12.5)):
    target_h = max_h
    target_w = int(target_h * view_w / view_h)
    if target_w > max_w:
        target_w = max_w
        target_h = int(target_w * view_h / view_w)
    x = int((SW - target_w) / 2)
    y = top
    slide.shapes.add_picture(path, x, y, width=target_w, height=target_h)

s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "System Architecture", "Three tiers. Full safety loop in under 100 ms.")
fit_image(s, f"{OUT}/fig_1_architecture.png", 900, 520)
add_footer(s, 5)


# ============ SLIDE 6 — Sensor Map ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "Sensor Placement on the Helmet", "Every component lives where its physics works best.")
fit_image(s, f"{OUT}/fig_2_sensor_map.png", 900, 560)
add_footer(s, 6)


# ============ SLIDE 7 — BOM ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "Bill of Materials", "Off-the-shelf parts. Prototype total ≈ ₹4,985. Retail at 10 k volume: ₹2,800.")

headers = ["#", "Component", "Model", "Placement", "Cost"]
bom = [
    ["1", "Microcontroller", "ESP32-S3-WROOM-1", "Rear pod", "₹450"],
    ["2", "6-axis IMU", "ICM-42688-P", "Crown", "₹380"],
    ["3", "LoRa radio", "Ra-02 (SX1278, 433 MHz)", "Rear pod + helical antenna", "₹220"],
    ["4", "GNSS", "u-blox NEO-6M", "Top shell, sky-facing", "₹350"],
    ["5", "Cellular modem", "SIM7600E-H 4G LTE", "Rear pod, chin-strap antenna", "₹1,800"],
    ["6", "Solar film", "Flexible a-Si 5V / 200 mA", "Top curved shell", "₹450"],
    ["7", "Battery", "Li-Po 3.7V 2000 mAh + BMS", "Rear pod, aramid sleeve", "₹350"],
    ["8", "MPPT IC", "CN3791", "Power board", "₹120"],
    ["9", "Buzzer / haptics / LEDs / bone-conduction", "Various", "Temple, ear-cup, brow", "₹400"],
    ["10", "Cancel button + QR + NFC + PCB", "—", "—", "₹465"],
]
table_x = Inches(0.5); table_y = Inches(1.7); table_w = Inches(12.3); table_h = Inches(4.8)
rows_count = len(bom) + 1
cols_w = [Inches(0.5), Inches(2.4), Inches(3.0), Inches(4.4), Inches(2.0)]
tbl = s.shapes.add_table(rows_count, 5, table_x, table_y, table_w, table_h).table
for i, w in enumerate(cols_w):
    tbl.columns[i].width = w

for i, h in enumerate(headers):
    cell = tbl.cell(0, i)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = h
    r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE; r.font.name='Calibri'
for ri, row in enumerate(bom, start=1):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = val
        r.font.size = Pt(11); r.font.color.rgb = NAVY; r.font.name='Calibri'
        if ri % 2 == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xf1, 0xf5, 0xf9)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE

# Total banner
add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5), AMBER, corner=True)
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         "TOTAL prototype BOM  ≈  ₹4,985    ·    Adds ~120 g to a stock helmet",
         size=15, bold=True, color=WHITE, align='center')
add_footer(s, 7)


# ============ SLIDE 8 — AI Pipeline ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "Edge-AI Pipeline — Two Models, One Chip", "PINN predicts the danger. TinyML confirms the crash. Both live on ESP32-S3.")
fit_image(s, f"{OUT}/fig_3_ai_pipeline.png", 900, 540)
add_footer(s, 8)


# ============ SLIDE 9 — Emergency Flow ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "Emergency Response Timeline", "Impact → ambulance dispatch in under 15 seconds. Every millisecond accounted for.")
fit_image(s, f"{OUT}/fig_4_emergency_flow.png", 900, 540)
add_footer(s, 9)


# ============ SLIDE 10 — LoRa Mesh ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "LoRa Rider Mesh", "One crash → every helmet in ~3 km knows in under 2 seconds. No cell tower needed.")
fit_image(s, f"{OUT}/fig_5_lora_mesh.png", 900, 540)
add_footer(s, 10)


# ============ SLIDE 11 — Solar ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "Solar Charging — Powered by the Sun", "Daytime rides leave the battery higher than they found it.")
fit_image(s, f"{OUT}/fig_6_solar_charging.png", 900, 560)
add_footer(s, 11)


# ============ SLIDE 12 — Why we win ============
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "Why RoadSoS Wins", "Six judging criteria — six clear answers.")
criteria = [
    ("Innovation", "First helmet to combine PINN + TinyML on a sub-₹500 MCU. Solar trickle is unique in this category.", AMBER),
    ("Feasibility", "All off-the-shelf parts. Working prototype in 6 weeks. BOM < ₹5,000.", TEAL),
    ("Impact", "Targets 44% of Indian road fatalities. 10% Golden Hour improvement → ~6,500 lives/year saved.", RED),
    ("Scalability", "LoRa mesh gets more valuable with each new helmet — classic network effect.", PURPLE),
    ("Sustainability", "Solar reduces charging burden, extends battery life, lowers e-waste.", TEAL),
    ("UX", "Zero rider input required. The helmet just works from the moment it's worn.", NAVY),
]
for i, (title, body, color) in enumerate(criteria):
    col, row = i % 2, i // 2
    x = Inches(0.5 + col * 6.25); y = Inches(1.85 + row * 1.55)
    add_rect(s, x, y, Inches(6.0), Inches(1.35), CARD, corner=True, line=color)
    add_rect(s, x, y, Inches(0.16), Inches(1.35), color)
    add_text(s, x + Inches(0.35), y + Inches(0.15), Inches(5.5), Inches(0.4), title, size=16, bold=True, color=NAVY)
    add_text(s, x + Inches(0.35), y + Inches(0.55), Inches(5.5), Inches(0.75), body, size=12, color=GREY)
add_footer(s, 12)


# ============ SLIDE 13 — Roadmap + Tagline ============
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY_DARK)
add_text(s, Inches(0.5), Inches(0.6), Inches(12), Inches(0.5),
         "ROADMAP & TAGLINE", size=14, bold=True, color=AMBER)
add_text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.8),
         "From hackathon to highway", size=36, bold=True, color=WHITE)

stages = [
    ("Wks 0–2", "Hackathon prototype"),
    ("Wks 3–6", "PINN trained on real data"),
    ("Month 3", "50-rider pilot (Dunzo / Swiggy)"),
    ("Month 6", "108 EMS + KA Police API"),
    ("Year 1", "B2C launch @ ₹4,999"),
]
for i, (t, b) in enumerate(stages):
    x = Inches(0.5 + i * 2.5)
    add_rect(s, x, Inches(2.5), Inches(2.3), Inches(1.7), RGBColor(0x15, 0x2e, 0x5e), corner=True)
    add_text(s, x, Inches(2.7), Inches(2.3), Inches(0.5), t, size=14, bold=True, color=AMBER, align='center')
    add_text(s, x, Inches(3.2), Inches(2.3), Inches(0.9), b, size=12, color=WHITE, align='center')

add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.8), RGBColor(0x15, 0x2e, 0x5e), corner=True)
add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
         "RoadSoS", size=32, bold=True, color=AMBER, align='center')
add_text(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.7),
         "Predicting the crash · Calling for help · Warning the world · Powered by the sun.",
         size=18, italic=True, color=WHITE, align='center')

add_text(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
         "Thank you. Questions?", size=12, color=LIGHT_GREY, align='center')


# Save
out_path = os.path.join(OUT, "RoadSoS_Deck.pptx")
prs.save(out_path)
print("Wrote:", out_path, "with", len(prs.slides), "slides")
